import streamlit as st
import asyncio
import aiohttp
import time
import re
from latin_cyrillic_symbols import to_cyrillic
from numberread import clean_text_number
from docx import Document
from pptx import Presentation
import pandas as pd
import fitz
import requests
import urllib3

api_key, API_URL, API_URL_HAND = st.secrets["API_TOKEN"], st.secrets["API_URL"],st.secrets["API_URL_HAND"]

def create_braille_symbols(symbol_type="alphabet"):
    symbols = {
        'alphabet': {
            'a': '⠁', 'b': '⠃', 'c': '⠉', 'd': '⠙', 'e': '⠑',
            'f': '⠋', 'g': '⠛', 'h': '⠓', 'i': '⠊', 'j': '⠚',
            'k': '⠅', 'l': '⠇', 'm': '⠍', 'n': '⠝', 'o': '⠕',
            'p': '⠏', 'q': '⠟', 'r': '⠗', 's': '⠎', 't': '⠞',
            'u': '⠥', 'v': '⠧', 'w': '⠺', 'x': '⠭', 'y': '⠽',
            'z': '⠵', 'ch': '⠟', 'sh': '⠱', "o'": '⠧', "g'": '⠻'
        },
        'numbers': {
            '0': '⠚', '1': '⠁', '2': '⠃', '3': '⠉', '4': '⠙',
            '5': '⠑', '6': '⠋', '7': '⠛', '8': '⠓', '9': '⠊'
        }
    }
    return symbols[symbol_type]

def clean_text(text):
    text = re.sub(r'(\d+)([a-zA-Z\u0400-\u04FF])', r'\1 \2', text)  
    text = re.sub(r'([a-zA-Z\u0400-\u04FF])(\d+)', r'\1 \2', text)
    text = re.sub(r'\s+', ' ', text)
    text = text.replace('\n', ' ')
    text = re.sub(r'[^a-zA-Z0-9\u0400-\u04FF\s]', '', text)
    return text.strip()

@st.cache_data(show_spinner=False)
def extract_text_from_docx(file):
    doc = Document(file)
    text = ' '.join(paragraph.text for paragraph in doc.paragraphs)
    return clean_text(text)

@st.cache_data(show_spinner=False)
def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ' '.join(paragraph.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame for paragraph in shape.text_frame.paragraphs)
    return clean_text(text)

@st.cache_data(show_spinner=False)
def extract_text_from_excel(file):
    excel_data = pd.read_excel(file, sheet_name=None)
    text = []
    for sheet_name, sheet_data in excel_data.items():
        text.append(sheet_data.to_string(index=False, header=True))
    return clean_text(' '.join(text))

@st.cache_data(show_spinner=False)
def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.read(), filetype='pdf')  
    text = []
    for page in doc:
        text.append(page.get_text("text"))
    return clean_text(' '.join(text))

@st.cache_data()
def handwrite_jpg_to_text(file_data):
    try:
        headers = {"Authorization": f"Bearer {st.secrets['API_TOKEN']}"}
        response = requests.post(st.secrets['API_URL_HAND'], headers=headers, data=file_data)
        if response.status_code == 200:
            return response.json()
    except Exception as e:
        st.warning(f"Xatolik bo'ldi: {e}")
        return None

#Audioga o'tkazish
def text_to_audio(text):    
    try:  
        headers = {"Content-Type": "application/json"}  
        data = {"text": text}  
        urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)  
        
        response = requests.post('https://oyqiz.airi.uz/api/v1/tts/', headers=headers, json=data, verify=False)  
        
        if response.status_code == 200:  
            return response.content 
        else:  
            return None   
    except Exception as ex:  
        st.error(f"xatolik: {ex}")  
        return None  
async def text_to_speech(text, language="uz"):
    try:
        text = text if any('\u0400' <= char <= '\u04FF' for char in text) else to_cyrillic(text)
        text = clean_text_number(text)
        headers = {"Authorization": f"Bearer {api_key}"}
        payload = {"inputs": text}
        async with aiohttp.ClientSession() as session:
            async with session.post(API_URL, headers=headers, json=payload) as response:
                if response.status == 200:
                    return await response.read()
                else:
                    error_message = await response.text()
                    st.error(f"⚠️ Xatolik: Server javobida xato (Status kodi: {response.status}). Tafsilotlar: {error_message}")
                    return None
    except Exception as e:
        st.toast("Internetga ulaning...", icon='❌')
        return None

async def main_text_to_speech():
    audio_content = text_to_audio(st.session_state.result)
    # audio_content = await text_to_speech(st.session_state.result, "uz")
    if audio_content:
        st.success("✅ Matn muvaffaqiyatli tovushga aylandi!")
        st.toast('Ajoyib!', icon='🎉')
        time.sleep(0.5)
        st.audio(audio_content, format='audio/wav', autoplay=True)
    else:
        st.warning("⚠️ Server bilan muammo yuz berdi. Qayta urinib ko'ring.")

def update_result():
    st.session_state.result = st.session_state.text_input

def app():
    if 'result' not in st.session_state:
        st.session_state.result = ""
    st.markdown("# ✅ :rainbow[Brayl klaviaturasi]")
    
    st.markdown("## 🔑 :red[Brayl Harflari]")
    alphabet = create_braille_symbols("alphabet")
    cols = st.columns(6)
    for i, (letter, symbol) in enumerate(alphabet.items()):
        with cols[i % 6]:
            if st.button(symbol, key=f"btn_{letter}", help=f"{letter}"):
                st.session_state.result += letter

    st.markdown("## 📌 :green[Brayl Raqamlari]")
    numbers = create_braille_symbols("numbers")
    num_cols = st.columns(5)
    for i, (number, symbol) in enumerate(numbers.items()):
        with num_cols[i % 5]:
            if st.button(symbol, key=f"num_btn_{number}", help=f"{number}"):
                st.session_state.result += number

    with st.sidebar:
        files = st.file_uploader("Faylingizni yuklang", type=['docx', 'pptx', 'xlsx', 'pdf', 'png'])
        if files:
            with st.spinner("Siz yuborgan fayldan matnlar ajratib olinayabdi. Iltimos, kuting..."):
                if files.name.endswith('.docx'):
                    st.session_state.result = extract_text_from_docx(files)
                elif files.name.endswith('.pptx'):
                    st.session_state.result = extract_text_from_pptx(files)
                elif files.name.endswith('.xlsx'):
                    st.session_state.result = extract_text_from_excel(files)
                elif files.name.endswith('.pdf'):
                    st.session_state.result = extract_text_from_pdf(files)
                elif files.name.endswith('.png'):
                    photo_file_bytes = files.read()
                    result_json = handwrite_jpg_to_text(photo_file_bytes)
                    if result_json:
                        st.session_state.result = clean_text(result_json[0].get('generated_text', ''))

        st.image('src/image.png', width=150)
        st.markdown("👁 :rainbow[O'zbekcha Brayl tarjimon]")

    row1, row2 = st.columns(2)
    with row1:
        if st.button("Bo'sh joy", key="space_button"):
            st.session_state.result += " "
    with row2:
        if st.button("Tozalash", key="clear_button"):
            st.session_state.result = ""
    
    st.text_area("💬 Matn kiriting yoki yuqoridagi Brayl klaviaturasidan foydalaning", 
                 value=st.session_state.result, 
                 key="text_input", 
                 height=100, 
                 max_chars=500, 
                 on_change=update_result)

    if st.button("🎧 Tovushga aylantirish", disabled=not st.session_state.result.strip()):
        with st.spinner("Jarayon boshlandi. Iltimos, kuting..."):
            asyncio.run(main_text_to_speech())

if __name__ == '__main__':
    app()
